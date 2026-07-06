"""
generation.py — PowerPoint de REX au style Square (template `data/template_REX.pptx`).

Une fiche = logo + titre 'DIRECTION – RÔLE' + résumé (IA mise en avant), placée
dans la grille 2×2 du layout "1_Références" du template (4 fiches/slide).
Le footer doré Square est hérité du master.
"""

from __future__ import annotations
import os
import re
import datetime
from collections import OrderedDict

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

from src.matching import Reference
from src.resume import resumer
from src.logos import chemin_logo

TEMPLATE = "data/template_REX.pptx"
LOGO_DIR = "data/logos"
FONT = "Verdana"
NOIR = RGBColor(0x1A, 0x1A, 0x1A)
GRIS = RGBColor(0x40, 0x40, 0x41)

# Ordre d'affichage des catégories (périmètres/types) dans le deck
PERIM_ORDER = ["Front Office", "Conformité", "Opérations",
               "Innovation métier", "Stratégie IT", "People & Change"]

# --- primitives de dessin ----------------------------------------------------
def _run(p, texte, taille, couleur, *, gras=False, ital=False):
    r = p.add_run(); r.text = texte
    f = r.font
    f.name = FONT; f.size = Pt(taille); f.bold = gras; f.italic = ital
    f.color.rgb = couleur
    return r


# --- Grille 2×2 du layout "1_Références" : 4 fiches/slide --------------------
# quadrants -> idx des placeholders (haut-gauche, haut-droite, bas-gauche, bas-droite)
QUADRANTS = [
    {"logo": 9,  "titre": 2, "desc": 1},
    {"logo": 14, "titre": 6, "desc": 5},
    {"logo": 13, "titre": 4, "desc": 3},
    {"logo": 15, "titre": 8, "desc": 7},
]


# --- Ajustement automatique du texte (aucun débordement de la boîte) ---------
# Hauteur des boîtes de description agrandies pour utiliser l'espace mort du
# template (l'écart entre la fin du texte et la rangée suivante) : la rangée
# du haut va jusqu'au séparateur (~3.93in), celle du bas jusqu'au footer
# (~6.80in). On garde une marge de sécurité.
DESC_H = [Inches(1.05), Inches(1.05), Inches(0.95), Inches(0.95)]  # par quadrant

# Largeur moyenne d'un caractère en Verdana, en pouces par point de police
# (calibrée empiriquement) ; majorée pour ne JAMAIS surestimer la place
# disponible (mieux vaut réduire un peu trop le texte que déborder).
_CAR_PAR_POUCE_PT = 0.0088


def _lignes_estimees(texte: str, taille: float, largeur_in: float) -> int:
    """Nombre de lignes que prendrait `texte` une fois justifié dans une boîte
    de largeur `largeur_in`, à la taille de police donnée (estimation)."""
    car_par_ligne = max(1, int(largeur_in / (taille * _CAR_PAR_POUCE_PT)))
    lignes, ligne = 1, ""
    for mot in texte.split():
        essai = f"{ligne} {mot}".strip()
        if len(essai) <= car_par_ligne or not ligne:
            ligne = essai
        else:
            lignes += 1
            ligne = mot
    return lignes


def _texte_ajuste(texte: str, largeur_in: float, hauteur_in: float, plafond_max: int):
    """Choisit un résumé + une taille de police qui tiennent SANS déborder
    dans une boîte largeur_in x hauteur_in (interligne 1.05)."""
    for taille in (9.5, 9.0, 8.5, 8.0, 7.5, 7.0):
        hauteur_ligne_in = taille * 1.2 * 1.05 / 72
        max_lignes = max(1, int(hauteur_in / hauteur_ligne_in))
        for max_car in range(plafond_max, 99, -20):
            court = resumer(texte, max_car=max_car)
            if _lignes_estimees(court, taille, largeur_in) <= max_lignes:
                return court, taille
    # dernier recours : texte très court à la plus petite taille
    return resumer(texte, max_car=100), 7.0


def _sans_tiret(txt: str) -> str:
    """Remplace les tirets longs (– —) par une virgule ; garde les traits
    d'union normaux (-)."""
    return (txt.replace(" – ", ", ").replace(" — ", ", ")
               .replace("–", "-").replace("—", "-"))


def _sans_retrait(p):
    """Supprime le retrait de 1re ligne / marge gauche hérités du layout."""
    pPr = p._p.get_or_add_pPr()
    pPr.set("indent", "0")
    pPr.set("marL", "0")


def _supprimer_ph(ph):
    """Retire réellement le placeholder (sinon son prompt 'Cliquer pour
    ajouter du contenu' reste visible en mode édition)."""
    ph._element.getparent().remove(ph._element)


def _titre_ph(ph, titre):
    tf = ph.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _sans_retrait(p)
    # taille adaptative pour éviter les titres sur 3 lignes qui frôlent le texte
    n = len(titre)
    taille = 11 if n <= 46 else 10 if n <= 70 else 9
    if " – " in titre:
        g, d = titre.split(" – ", 1)
        _run(p, _sans_tiret(g).upper() + ", ", taille, NOIR, gras=True)
        _run(p, _sans_tiret(d).upper(), taille, NOIR, ital=True)
    else:
        _run(p, _sans_tiret(titre).upper(), taille, NOIR, gras=True)


def _desc_ph(ph, texte, taille=9.5):
    tf = ph.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.JUSTIFY; p.line_spacing = 1.05
    _sans_retrait(p)
    _run(p, texte, taille, GRIS)


def _logo_ph(slide, ph, path):
    L, T, W, H = ph.left, ph.top, ph.width, ph.height
    with Image.open(path) as im:
        iw, ih = im.size
    ratio = min(W / iw, H / ih)
    w, h = int(iw * ratio), int(ih * ratio)
    slide.shapes.add_picture(path, L + (W - w) // 2, T + (H - h) // 2, w, h)
    _supprimer_ph(ph)          # on retire le placeholder vide (prompt)


def _nom_ph(ph, nom):
    """Fallback quand aucun logo n'est disponible : on affiche le nom du client
    (centré, gras) dans la zone logo."""
    tf = ph.text_frame; tf.clear(); tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _sans_retrait(p)
    _run(p, nom, 14, NOIR, gras=True)


def _slide_4up(prs, layout, titre, lot):
    s = prs.slides.add_slide(layout)
    s.placeholders[0].text = titre
    for i, q in enumerate(QUADRANTS):
        if i < len(lot):
            ref, logo = lot[i]
            _titre_ph(s.placeholders[q["titre"]], ref.titre or ref.perimetre)
            # boîte de description agrandie (utilise l'espace mort du template).
            # Attention : (re)fixer left/top/width en même temps que height,
            # sinon python-pptx remet left/top/width à 0 (position héritée du
            # layout perdue dès qu'on écrit UN SEUL attribut de xfrm).
            desc_ph = s.placeholders[q["desc"]]
            l, t, w = desc_ph.left, desc_ph.top, desc_ph.width
            desc_ph.left, desc_ph.top, desc_ph.width, desc_ph.height = l, t, w, DESC_H[i]
            largeur_in, hauteur_in = desc_ph.width / 914400, desc_ph.height / 914400
            plafond = 420 if i < 2 else 320    # rangée du bas : résumé plus court
            texte, taille = _texte_ajuste(ref.description, largeur_in, hauteur_in, plafond)
            _desc_ph(desc_ph, texte, taille)
            if logo:
                _logo_ph(s, s.placeholders[q["logo"]], logo)
            else:
                _nom_ph(s.placeholders[q["logo"]], ref.client)
        else:  # quadrant inutilisé -> on retire tous ses placeholders
            for k in ("titre", "desc", "logo"):
                _supprimer_ph(s.placeholders[q[k]])
    return s


_DATE_RE = re.compile(r"^\d{2}/\d{2}/\d{4}$")


def _corriger_date(prs):
    """Le template contient une date figée (12/05/2023) en texte dans le master.
    On la remplace par la date du jour (footer de toutes les slides)."""
    auj = datetime.date.today().strftime("%d/%m/%Y")

    def _fix(shapes):
        for sh in shapes:
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if _DATE_RE.match(r.text.strip()):
                            r.text = auj

    for master in prs.slide_masters:
        _fix(master.shapes)
        for layout in master.slide_layouts:
            _fix(layout.shapes)


# --- Point d'entrée ----------------------------------------------------------
def construire_deck(refs: list[Reference], sortie: str) -> str:
    """4 fiches/slide, texte raccourci (grille 2×2 du template)."""
    prs = Presentation(TEMPLATE)
    lst = prs.slides._sldIdLst
    for sldId in list(lst):
        rId = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        if rId:
            prs.part.drop_rel(rId)
        lst.remove(sldId)

    _corriger_date(prs)                       # footer : date du jour au lieu de 2023

    layout = prs.slide_layouts[1]

    # regroupement par CATÉGORIE (périmètre), ordre de pertinence conservé
    # dans chaque groupe ; les cartes d'une slide peuvent venir de banques
    # différentes (chacune garde son logo).
    groupes: "OrderedDict[str, list]" = OrderedDict()
    for perim in PERIM_ORDER:                 # ordre fixe des catégories
        for r in refs:
            if r.perimetre == perim:
                groupes.setdefault(perim, []).append((r, chemin_logo(r.client)))

    for perim, lot in groupes.items():
        pages = [lot[i:i + 4] for i in range(0, len(lot), 4)]
        for idx, page in enumerate(pages, 1):
            titre = f"Références IA : {perim}"
            if len(pages) > 1:
                titre += f" ({idx}/{len(pages)})"
            _slide_4up(prs, layout, titre, page)

    os.makedirs(os.path.dirname(sortie), exist_ok=True)
    prs.save(sortie)
    return sortie
